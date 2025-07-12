from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
import os
import shutil
from pathlib import Path
import pypdf
from docx import Document
from openpyxl import Workbook
import tempfile
import zipfile
from typing import List
import uuid

app = FastAPI(title="PDF Platform", description="Plataforma de manipulação de PDFs")

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Diretórios
UPLOAD_DIR = "uploads"
OUTPUT_DIR = "output"
Path(UPLOAD_DIR).mkdir(exist_ok=True)
Path(OUTPUT_DIR).mkdir(exist_ok=True)

app.mount("/static", StaticFiles(directory="static"), name="static")

@app.get("/", response_class=HTMLResponse)
async def get_frontend():
    return """
    <!DOCTYPE html>
    <html lang="pt-BR">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>PDF Platform - Lobios</title>
        <link rel="icon" type="image/png" href="/static/logo.png"/>
        <link href="https://fonts.googleapis.com/css2?family=Quicksand:wght@500;700&display=swap" rel="stylesheet">
        <style>
            :root {
                --lobios-purple: #7b3294;
                --lobios-purple-light: #a259c6;
                --lobios-bg: #f8f8fa;
                --lobios-card: #fff;
                --lobios-gray: #e5e5e5;
                --lobios-dark: #222;
            }
            body {
                margin: 0; padding: 0; font-family: 'Segoe UI', Arial, sans-serif;
                background: var(--lobios-bg);
                color: var(--lobios-dark);
            }
            .sidebar {
                position: fixed; left: 0; top: 0; bottom: 0; width: 220px; background: #fff;
                color: var(--lobios-purple); display: flex; flex-direction: column; align-items: center; padding-top: 30px; z-index: 2; border-right: 1.5px solid #eee;
            }
            .sidebar img {
                width: 140px; margin-bottom: 30px; filter: none;
            }
            .sidebar nav {
                width: 100%;
            }
            .sidebar nav a {
                display: block; color: var(--lobios-purple); text-decoration: none; padding: 14px 30px; font-size: 16px;
                border-left: 4px solid transparent; transition: background 0.2s, border 0.2s; font-weight: 500;
            }
            .sidebar nav a.active, .sidebar nav a:hover {
                background: #f3eafd; border-left: 4px solid var(--lobios-purple);
                color: var(--lobios-purple);
            }
            .main {
                margin-left: 220px; min-height: 100vh;
            }
            .topbar {
                background: #fff; height: 64px; display: flex; align-items: center; justify-content: flex-end;
                box-shadow: 0 2px 8px rgba(123,50,148,0.07); padding: 0 40px; position: sticky; top: 0; z-index: 1;
            }
            .topbar .user {
                font-weight: 500; color: var(--lobios-purple); display: flex; align-items: center;
            }
            .topbar .user:before {
                content: '\1F464'; font-size: 22px; margin-right: 8px;
            }
            .topbar .impact {
                font-family: 'Quicksand', 'Montserrat', 'Segoe UI', Arial, sans-serif;
                font-weight: 700;
                color: var(--lobios-purple);
                font-size: 1.35rem;
                letter-spacing: 0.01em;
                text-align: right;
                width: 100%;
                margin-top: 2px;
            }
            .container {
                max-width: 1200px; margin: 30px auto; padding: 0 20px;
            }
            .header h1 {
                color: var(--lobios-purple); margin-bottom: 10px; font-size: 2.2rem;
            }
            .tools-grid {
                display: grid; grid-template-columns: repeat(auto-fit, minmax(320px, 1fr)); gap: 24px;
            }
            .tool-card {
                background: var(--lobios-card); border-radius: 12px; padding: 28px 22px; box-shadow: 0 2px 12px rgba(123,50,148,0.08);
                display: flex; flex-direction: column; align-items: stretch;
            }
            .tool-card h3 {
                color: var(--lobios-purple); margin-bottom: 18px; font-size: 19px; font-weight: 600;
            }
            .file-input { width: 100%; padding: 10px; border: 2px dashed var(--lobios-purple-light); border-radius: 6px; margin-bottom: 15px; cursor: pointer; background: #faf7fc; }
            .file-input:hover { border-color: var(--lobios-purple); }
            .btn { background: var(--lobios-purple); color: white; border: none; padding: 13px 0; border-radius: 6px; cursor: pointer; width: 100%; font-size: 15px; font-weight: 500; transition: background 0.2s; }
            .btn:hover { background: var(--lobios-purple-light); }
            .btn:disabled { background: #bdc3c7; cursor: not-allowed; }
            .result { margin-top: 15px; padding: 10px; background: #e6e6fa; border-radius: 6px; display: none; color: var(--lobios-dark); }
            .error { background: #f8d7da; color: #721c24; }
            .loading { display: none; text-align: center; margin-top: 10px; }
            input[type="number"] { width: 100%; padding: 8px; margin: 5px 0; border: 1px solid #ddd; border-radius: 4px; }
            @media (max-width: 900px) {
                .sidebar { width: 60px; padding-top: 18px; }
                .sidebar img { width: 38px; margin-bottom: 18px; }
                .sidebar nav a { font-size: 0; padding: 12px 10px; }
                .main { margin-left: 60px; }
            }
        </style>
    </head>
    <body>
        <div class="sidebar">
            <img src="/static/logo.png" alt="Lobios">
            <nav id="sidebar-categories">
                <a href="#" class="category-link active" data-category="organizar">Organizar PDF</a>
                <a href="#" class="category-link" data-category="otimizar">Otimizar PDF</a>
                <a href="#" class="category-link" data-category="converter-em">Converter em PDF</a>
                <a href="#" class="category-link" data-category="converter-de">Converter de PDF</a>
                <a href="#" class="category-link" data-category="editar">Editar PDF</a>
                <a href="#" class="category-link" data-category="seguranca">Segurança do PDF</a>
            </nav>
        </div>
        <div class="main">
            <div class="topbar">
                <div class="impact">Ferramentas online para facilitar o manuseio de arquivos PDF com segurança e praticidade.</div>
            </div>
            <div class="container">
                <div class="header">
                    <h1 id="category-title">Organizar PDF</h1>
                    <p id="category-desc">Manipule seus arquivos PDF facilmente</p>
                </div>
                <div id="category-functions">
                    <!-- ORGANIZAR PDF -->
                    <div class="tools-grid" id="organizar" style="display: grid;">
                        <!-- Juntar PDF (funcional) -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">🧩 Juntar PDF</h3>
                            <input type="file" id="mergePdfs" accept=".pdf" multiple class="file-input input-full">
                            <button onclick="mergePdfs()" class="btn input-full">Juntar Arquivos</button>
                            <div class="loading" id="loadingMerge">Juntando...</div>
                            <div class="result" id="resultMerge"></div>
                        </div>
                        <!-- Dividir PDF (funcional) -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">✂️ Dividir PDF</h3>
                            <input type="file" id="splitPdf" accept=".pdf" class="file-input input-full">
                            <input type="number" placeholder="Página inicial" id="startPage" min="1" class="input-full">
                            <input type="number" placeholder="Página final" id="endPage" min="1" class="input-full">
                            <button onclick="splitPdf()" class="btn input-full">Extrair Páginas</button>
                            <div class="loading" id="loadingSplit">Extraindo...</div>
                            <div class="result" id="resultSplit"></div>
                        </div>
                        <!-- Remover páginas (nova função) -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">❌ Remover páginas</h3>
                            <input type="file" id="removePagesPdf" accept=".pdf" class="file-input input-full">
                            <input type="text" placeholder="Ex: 2,4,7-9" id="removePagesList" class="input-full">
                            <button onclick="removePages()" class="btn input-full">Remover Páginas</button>
                            <div class="loading" id="loadingRemovePages">Removendo...</div>
                            <div class="result" id="resultRemovePages"></div>
                        </div>
                        <!-- Extrair páginas (nova função) -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">📤 Extrair páginas</h3>
                            <input type="file" id="extractPagesPdf" accept=".pdf" class="file-input input-full">
                            <input type="text" placeholder="Ex: 1,3,5-7" id="extractPagesList" class="input-full">
                            <button onclick="extractPages()" class="btn input-full">Extrair Páginas</button>
                            <div class="loading" id="loadingExtractPages">Extraindo...</div>
                            <div class="result" id="resultExtractPages"></div>
                        </div>
                        <!-- Organizar PDF (nova função) -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">🔀 Organizar PDF</h3>
                            <input type="file" id="organizePdf" accept=".pdf" class="file-input input-full">
                            <input type="text" placeholder="Nova ordem (Ex: 3,1,2,5,4)" id="organizeOrder" class="input-full">
                            <button onclick="organizePdfPages()" class="btn input-full">Organizar Páginas</button>
                            <div class="loading" id="loadingOrganizePdf">Organizando...</div>
                            <div class="result" id="resultOrganizePdf"></div>
                        </div>
                        <!-- Digitalizar PDF (em breve) -->
                        <div class="tool-card disabled"><h3>📷 Digitalizar PDF</h3><p>Em breve</p></div>
                    </div>
                    <!-- OTIMIZAR PDF -->
                    <div class="tools-grid" id="otimizar" style="display: none;">
                        <!-- Comprimir PDF (funcional) -->
                        <div class="tool-card">
                            <h3>🗜️ Comprimir PDF</h3>
                            <input type="file" id="compressPdf" accept=".pdf" class="file-input">
                            <button onclick="compressPdf()" class="btn">Comprimir Arquivo</button>
                            <div class="loading" id="loadingCompress">Comprimindo...</div>
                            <div class="result" id="resultCompress"></div>
                        </div>
                        <!-- Em breve -->
                        <div class="tool-card disabled"><h3>🛠️ Reparar PDF</h3><p>Em breve</p></div>
                        <div class="tool-card disabled"><h3>📝 OCR PDF</h3><p>Em breve</p></div>
                    </div>
                    <!-- CONVERTER EM PDF -->
                    <div class="tools-grid" id="converter-em" style="display: none;">
                        <!-- JPG para PDF -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">🖼️ JPG para PDF</h3>
                            <input type="file" id="jpgToPdf" accept=".jpg,.jpeg,.png" multiple class="file-input input-full">
                            <button onclick="convertJpgToPdf()" class="btn input-full">Converter para PDF</button>
                            <div class="loading" id="loadingJpgToPdf">Convertendo...</div>
                            <div class="result" id="resultJpgToPdf"></div>
                        </div>
                        <!-- WORD para PDF -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">📝 WORD para PDF</h3>
                            <input type="file" id="wordToPdf" accept=".doc,.docx" class="file-input input-full">
                            <button onclick="convertWordToPdf()" class="btn input-full">Converter para PDF</button>
                            <div class="loading" id="loadingWordToPdf">Convertendo...</div>
                            <div class="result" id="resultWordToPdf"></div>
                        </div>
                        <!-- EXCEL para PDF -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">📊 EXCEL para PDF</h3>
                            <input type="file" id="excelToPdf" accept=".xls,.xlsx" class="file-input input-full">
                            <button onclick="convertExcelToPdf()" class="btn input-full">Converter para PDF</button>
                            <div class="loading" id="loadingExcelToPdf">Convertendo...</div>
                            <div class="result" id="resultExcelToPdf"></div>
                        </div>
                        <!-- POWERPOINT para PDF -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">📈 POWERPOINT para PDF</h3>
                            <input type="file" id="pptToPdf" accept=".ppt,.pptx" class="file-input input-full">
                            <button onclick="convertPptToPdf()" class="btn input-full">Converter para PDF</button>
                            <div class="loading" id="loadingPptToPdf">Convertendo...</div>
                            <div class="result" id="resultPptToPdf"></div>
                        </div>
                        <!-- HTML para PDF -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">🌐 HTML para PDF</h3>
                            <input type="file" id="htmlToPdf" accept=".html,.htm" class="file-input input-full">
                            <button onclick="convertHtmlToPdf()" class="btn input-full">Converter para PDF</button>
                            <div class="loading" id="loadingHtmlToPdf">Convertendo...</div>
                            <div class="result" id="resultHtmlToPdf"></div>
                        </div>
                    </div>
                    <!-- CONVERTER DE PDF -->
                    <div class="tools-grid" id="converter-de" style="display: none;">
                        <div class="tool-card disabled"><h3>🖼️ PDF para JPG</h3><p>Em breve</p></div>
                        <div class="tool-card">
                            <h3>📝 PDF para WORD</h3>
                            <input type="file" id="pdfToWord" accept=".pdf" class="file-input">
                            <button onclick="convertToWord()" class="btn">Converter para DOCX</button>
                            <div class="loading" id="loadingWord">Convertendo...</div>
                            <div class="result" id="resultWord"></div>
                        </div>
                        <div class="tool-card">
                            <h3>📊 PDF para EXCEL</h3>
                            <input type="file" id="pdfToExcel" accept=".pdf" class="file-input">
                            <button onclick="convertToExcel()" class="btn">Converter para XLSX</button>
                            <div class="loading" id="loadingExcel">Convertendo...</div>
                            <div class="result" id="resultExcel"></div>
                        </div>
                        <div class="tool-card disabled"><h3>📈 PDF para POWERPOINT</h3><p>Em breve</p></div>
                        <div class="tool-card disabled"><h3>🅰️ PDF para PDF/A</h3><p>Em breve</p></div>
                    </div>
                    <!-- EDITAR PDF -->
                    <div class="tools-grid" id="editar" style="display: none;">
                        <!-- Rodar PDF (em breve) -->
                        <div class="tool-card disabled"><h3>🔄 Rodar PDF</h3><p>Em breve</p></div>
                        <!-- Inserir números de página (funcional) -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">#️⃣ Inserir números de página</h3>
                            <input type="file" id="numberPdf" accept=".pdf" class="file-input input-full">
                            <button onclick="insertPageNumbers()" class="btn input-full">Inserir Números</button>
                            <div class="loading" id="loadingNumberPdf">Processando...</div>
                            <div class="result" id="resultNumberPdf"></div>
                        </div>
                        <!-- Inserir marca d'água (funcional) -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">💧 Inserir marca d'água</h3>
                            <input type="file" id="watermarkPdf" accept=".pdf" class="file-input input-full">
                            <input type="text" id="watermarkText" placeholder="Texto da marca d'água" class="input-full">
                            <button onclick="insertWatermark()" class="btn input-full">Inserir Marca d'Água</button>
                            <div class="loading" id="loadingWatermarkPdf">Processando...</div>
                            <div class="result" id="resultWatermarkPdf"></div>
                        </div>
                        <!-- Recortar PDF (em breve) -->
                        <div class="tool-card disabled"><h3>✂️ Recortar PDF</h3><p>Em breve</p></div>
                        <!-- Editar PDF (em breve) -->
                        <div class="tool-card disabled"><h3>✏️ Editar PDF</h3><p>Em breve</p></div>
                    </div>
                    <!-- SEGURANÇA DO PDF -->
                    <div class="tools-grid" id="seguranca" style="display: none;">
                        <!-- Desbloquear PDF -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">🔓 Desbloquear PDF</h3>
                            <input type="file" id="unlockPdf" accept=".pdf" class="file-input input-full">
                            <input type="password" placeholder="Senha atual" id="unlockPassword" class="input-full">
                            <button onclick="unlockPdf()" class="btn input-full">Desbloquear PDF</button>
                            <div class="loading" id="loadingUnlockPdf">Desbloqueando...</div>
                            <div class="result" id="resultUnlockPdf"></div>
                        </div>
                        <!-- Proteger PDF -->
                        <div class="tool-card">
                            <h3 style="color: var(--lobios-purple)">🛡️ Proteger PDF</h3>
                            <input type="file" id="protectPdf" accept=".pdf" class="file-input input-full">
                            <input type="password" placeholder="Nova senha" id="protectPassword" class="input-full">
                            <button onclick="protectPdf()" class="btn input-full">Proteger PDF</button>
                            <div class="loading" id="loadingProtectPdf">Protegendo...</div>
                            <div class="result" id="resultProtectPdf"></div>
                        </div>
                        <!-- Outras funções -->
                        <div class="tool-card disabled"><h3>🖊️ Assinar PDF</h3><p>Em breve</p></div>
                        <div class="tool-card disabled"><h3>🙈 Ocultar PDF</h3><p>Em breve</p></div>
                        <div class="tool-card">
                            <h3>📋 Comparar PDF</h3>
                            <input type="file" id="comparePdf1" accept=".pdf" class="file-input input-full" placeholder="PDF 1">
                            <input type="file" id="comparePdf2" accept=".pdf" class="file-input input-full" placeholder="PDF 2">
                            <button onclick="comparePdfs()" class="btn input-full">Comparar Arquivos</button>
                            <div class="loading" id="loadingCompare">Comparando...</div>
                            <div class="result" id="resultCompare"></div>
                        </div>
                    </div>
                </div>
            </div>
        </div>
        <!-- Botão Política de Privacidade -->
        <button id="privacyBtn" style="position:fixed;right:24px;bottom:24px;z-index:99;background:#fff;color:var(--lobios-purple);border:1.5px solid #eee;padding:10px 22px;border-radius:24px;font-weight:600;box-shadow:0 2px 8px rgba(123,50,148,0.07);cursor:pointer;">Política de Privacidade</button>
        <!-- Modal Política -->
        <div id="privacyModal" style="display:none;position:fixed;top:0;left:0;width:100vw;height:100vh;background:rgba(0,0,0,0.18);z-index:100;align-items:center;justify-content:center;">
            <div style="background:#fff;border-radius:18px;max-width:900px;width:95vw;padding:36px 28px;box-shadow:0 4px 32px rgba(0,0,0,0.13);position:relative;">
                <span id="closePrivacy" style="position:absolute;top:18px;right:24px;font-size:28px;cursor:pointer;color:#aaa;">&times;</span>
                <h2 style="color:var(--lobios-purple);margin-bottom:8px;">Políticas de segurança e privacidade de dados</h2>
                <p style="margin-bottom:28px;color:#444;">Informações detalhadas sobre a estrutura de privacidade e segurança do PDF Platform Lobios.</p>
                <div style="display:grid;grid-template-columns:repeat(auto-fit,minmax(200px,1fr));gap:22px;">
                    <div style="background:#faf7fc;border-radius:12px;padding:22px 16px;text-align:center;">
                        <div style="font-size:36px;">🔒</div>
                        <h4 style="color:var(--lobios-purple);margin:10px 0 8px 0;">Segurança</h4>
                        <p style="font-size:15px;color:#333;">Todos os arquivos enviados são processados de forma segura e excluídos automaticamente após a conversão. Não armazenamos, visualizamos ou compartilhamos seus documentos.</p>
                    </div>
                    <div style="background:#faf7fc;border-radius:12px;padding:22px 16px;text-align:center;">
                        <div style="font-size:36px;">🛡️</div>
                        <h4 style="color:var(--lobios-purple);margin:10px 0 8px 0;">Privacidade</h4>
                        <p style="font-size:15px;color:#333;">Sua privacidade é prioridade. Os arquivos são eliminados dos nossos servidores logo após o processamento, garantindo total confidencialidade.</p>
                    </div>
                    <div style="background:#faf7fc;border-radius:12px;padding:22px 16px;text-align:center;">
                        <div style="font-size:36px;">📄</div>
                        <h4 style="color:var(--lobios-purple);margin:10px 0 8px 0;">Termos</h4>
                        <p style="font-size:15px;color:#333;">Ao utilizar o PDF Platform Lobios, você concorda com nossos termos: não armazenamos arquivos, não compartilhamos dados e não utilizamos seus documentos para nenhum outro fim.</p>
                    </div>
                    <div style="background:#faf7fc;border-radius:12px;padding:22px 16px;text-align:center;">
                        <div style="font-size:36px;">🍪</div>
                        <h4 style="color:var(--lobios-purple);margin:10px 0 8px 0;">Cookies</h4>
                        <p style="font-size:15px;color:#333;">Utilizamos apenas cookies essenciais para o funcionamento da plataforma. Não rastreamos, não vendemos e não utilizamos cookies para fins de marketing.</p>
                    </div>
                </div>
            </div>
        </div>
        <script>
            const API_BASE = '';

            function showLoading(id) {
                document.getElementById(`loading${id}`).style.display = 'block';
                document.getElementById(`result${id}`).style.display = 'none';
            }

            function hideLoading(id) {
                document.getElementById(`loading${id}`).style.display = 'none';
            }

            function showResult(id, message, isError = false) {
                const result = document.getElementById(`result${id}`);
                result.innerHTML = message;
                result.className = isError ? 'result error' : 'result';
                result.style.display = 'block';
            }

            async function convertToWord() {
                const file = document.getElementById('pdfToWord').files[0];
                if (!file) return alert('Selecione um arquivo PDF');

                showLoading('Word');
                const formData = new FormData();
                formData.append('file', file);

                try {
                    const response = await fetch('/convert/word', {
                        method: 'POST',
                        body: formData
                    });

                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        showResult('Word', `<a href="${url}" download="${file.name.replace('.pdf', '.docx')}">📥 Download DOCX</a>`);
                    } else {
                        throw new Error('Erro na conversão');
                    }
                } catch (error) {
                    showResult('Word', 'Erro ao converter arquivo', true);
                } finally {
                    hideLoading('Word');
                }
            }

            async function convertToExcel() {
                const file = document.getElementById('pdfToExcel').files[0];
                if (!file) return alert('Selecione um arquivo PDF');

                showLoading('Excel');
                const formData = new FormData();
                formData.append('file', file);

                try {
                    const response = await fetch('/convert/excel', {
                        method: 'POST',
                        body: formData
                    });

                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        showResult('Excel', `<a href="${url}" download="${file.name.replace('.pdf', '.xlsx')}">📥 Download XLSX</a>`);
                    } else {
                        throw new Error('Erro na conversão');
                    }
                } catch (error) {
                    showResult('Excel', 'Erro ao converter arquivo', true);
                } finally {
                    hideLoading('Excel');
                }
            }

            async function mergePdfs() {
                const files = document.getElementById('mergePdfs').files;
                if (files.length < 2) return alert('Selecione pelo menos 2 arquivos PDF');

                showLoading('Merge');
                const formData = new FormData();
                for (let file of files) {
                    formData.append('files', file);
                }

                try {
                    const response = await fetch('/merge', {
                        method: 'POST',
                        body: formData
                    });

                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        showResult('Merge', `<a href="${url}" download="merged.pdf">📥 Download PDF Combinado</a>`);
                    } else {
                        throw new Error('Erro ao juntar arquivos');
                    }
                } catch (error) {
                    showResult('Merge', 'Erro ao juntar arquivos', true);
                } finally {
                    hideLoading('Merge');
                }
            }

            async function splitPdf() {
                const file = document.getElementById('splitPdf').files[0];
                const startPage = document.getElementById('startPage').value;
                const endPage = document.getElementById('endPage').value;

                if (!file) return alert('Selecione um arquivo PDF');
                if (!startPage || !endPage) return alert('Informe as páginas inicial e final');

                showLoading('Split');
                const formData = new FormData();
                formData.append('file', file);
                formData.append('start_page', startPage);
                formData.append('end_page', endPage);

                try {
                    const response = await fetch('/split', {
                        method: 'POST',
                        body: formData
                    });

                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        showResult('Split', `<a href="${url}" download="extracted_pages.pdf">📥 Download Páginas Extraídas</a>`);
                    } else {
                        throw new Error('Erro ao extrair páginas');
                    }
                } catch (error) {
                    showResult('Split', 'Erro ao extrair páginas', true);
                } finally {
                    hideLoading('Split');
                }
            }

            async function compressPdf() {
                const file = document.getElementById('compressPdf').files[0];
                if (!file) return alert('Selecione um arquivo PDF');

                showLoading('Compress');
                const formData = new FormData();
                formData.append('file', file);

                try {
                    const response = await fetch('/compress', {
                        method: 'POST',
                        body: formData
                    });

                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        showResult('Compress', `<a href="${url}" download="${file.name.replace('.pdf', '_compressed.pdf')}">📥 Download PDF Comprimido</a>`);
                    } else {
                        throw new Error('Erro ao comprimir');
                    }
                } catch (error) {
                    showResult('Compress', 'Erro ao comprimir arquivo', true);
                } finally {
                    hideLoading('Compress');
                }
            }

            async function comparePdfs() {
                const file1 = document.getElementById('comparePdf1').files[0];
                const file2 = document.getElementById('comparePdf2').files[0];

                if (!file1 || !file2) return alert('Selecione os dois arquivos PDF');

                showLoading('Compare');
                const formData = new FormData();
                formData.append('file1', file1);
                formData.append('file2', file2);

                try {
                    const response = await fetch('/compare', {
                        method: 'POST',
                        body: formData
                    });

                    const result = await response.json();
                    if (response.ok) {
                        showResult('Compare', result.message);
                    } else {
                        throw new Error(result.detail);
                    }
                } catch (error) {
                    showResult('Compare', 'Erro ao comparar arquivos', true);
                } finally {
                    hideLoading('Compare');
                }
            }

            async function removePages() {
                const file = document.getElementById('removePagesPdf').files[0];
                const pages = document.getElementById('removePagesList').value;
                if (!file || !pages) return alert('Selecione o PDF e informe as páginas a remover');
                document.getElementById('loadingRemovePages').style.display = 'block';
                document.getElementById('resultRemovePages').style.display = 'none';
                const formData = new FormData();
                formData.append('file', file);
                formData.append('pages', pages);
                try {
                    const response = await fetch('/remove-pages', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultRemovePages').innerHTML = `<a href="${url}" download="removed_pages.pdf">📥 Download PDF sem páginas</a>`;
                        document.getElementById('resultRemovePages').className = 'result';
                    } else {
                        throw new Error('Erro ao remover páginas');
                    }
                } catch (error) {
                    document.getElementById('resultRemovePages').innerHTML = 'Erro ao remover páginas';
                    document.getElementById('resultRemovePages').className = 'result error';
                } finally {
                    document.getElementById('loadingRemovePages').style.display = 'none';
                    document.getElementById('resultRemovePages').style.display = 'block';
                }
            }
            async function extractPages() {
                const file = document.getElementById('extractPagesPdf').files[0];
                const pages = document.getElementById('extractPagesList').value;
                if (!file || !pages) return alert('Selecione o PDF e informe as páginas a extrair');
                document.getElementById('loadingExtractPages').style.display = 'block';
                document.getElementById('resultExtractPages').style.display = 'none';
                const formData = new FormData();
                formData.append('file', file);
                formData.append('pages', pages);
                try {
                    const response = await fetch('/extract-pages', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultExtractPages').innerHTML = `<a href="${url}" download="extracted_pages.pdf">📥 Download Páginas Extraídas</a>`;
                        document.getElementById('resultExtractPages').className = 'result';
                    } else {
                        throw new Error('Erro ao extrair páginas');
                    }
                } catch (error) {
                    document.getElementById('resultExtractPages').innerHTML = 'Erro ao extrair páginas';
                    document.getElementById('resultExtractPages').className = 'result error';
                } finally {
                    document.getElementById('loadingExtractPages').style.display = 'none';
                    document.getElementById('resultExtractPages').style.display = 'block';
                }
            }
            async function organizePdfPages() {
                const file = document.getElementById('organizePdf').files[0];
                const order = document.getElementById('organizeOrder').value;
                if (!file || !order) return alert('Selecione o PDF e informe a nova ordem das páginas');
                document.getElementById('loadingOrganizePdf').style.display = 'block';
                document.getElementById('resultOrganizePdf').style.display = 'none';
                const formData = new FormData();
                formData.append('file', file);
                formData.append('order', order);
                try {
                    const response = await fetch('/organize-pages', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultOrganizePdf').innerHTML = `<a href="${url}" download="organized.pdf">📥 Download PDF Organizado</a>`;
                        document.getElementById('resultOrganizePdf').className = 'result';
                    } else {
                        throw new Error('Erro ao organizar páginas');
                    }
                } catch (error) {
                    document.getElementById('resultOrganizePdf').innerHTML = 'Erro ao organizar páginas';
                    document.getElementById('resultOrganizePdf').className = 'result error';
                } finally {
                    document.getElementById('loadingOrganizePdf').style.display = 'none';
                    document.getElementById('resultOrganizePdf').style.display = 'block';
                }
            }

            async function convertJpgToPdf() {
                const files = document.getElementById('jpgToPdf').files;
                if (!files.length) return alert('Selecione pelo menos uma imagem');
                document.getElementById('loadingJpgToPdf').style.display = 'block';
                document.getElementById('resultJpgToPdf').style.display = 'none';
                const formData = new FormData();
                for (let file of files) formData.append('files', file);
                try {
                    const response = await fetch('/convert/jpg-to-pdf', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultJpgToPdf').innerHTML = `<a href="${url}" download="imagens.pdf">📥 Download PDF</a>`;
                        document.getElementById('resultJpgToPdf').className = 'result';
                    } else {
                        throw new Error('Erro na conversão');
                    }
                } catch (error) {
                    document.getElementById('resultJpgToPdf').innerHTML = 'Erro ao converter';
                    document.getElementById('resultJpgToPdf').className = 'result error';
                } finally {
                    document.getElementById('loadingJpgToPdf').style.display = 'none';
                    document.getElementById('resultJpgToPdf').style.display = 'block';
                }
            }
            async function convertWordToPdf() {
                const file = document.getElementById('wordToPdf').files[0];
                if (!file) return alert('Selecione um arquivo Word');
                document.getElementById('loadingWordToPdf').style.display = 'block';
                document.getElementById('resultWordToPdf').style.display = 'none';
                const formData = new FormData();
                formData.append('file', file);
                try {
                    const response = await fetch('/convert/word-to-pdf', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultWordToPdf').innerHTML = `<a href="${url}" download="${file.name.replace(/\.[^.]+$/, '.pdf')}">📥 Download PDF</a>`;
                        document.getElementById('resultWordToPdf').className = 'result';
                    } else {
                        throw new Error('Erro na conversão');
                    }
                } catch (error) {
                    document.getElementById('resultWordToPdf').innerHTML = 'Erro ao converter';
                    document.getElementById('resultWordToPdf').className = 'result error';
                } finally {
                    document.getElementById('loadingWordToPdf').style.display = 'none';
                    document.getElementById('resultWordToPdf').style.display = 'block';
                }
            }
            async function convertExcelToPdf() {
                const file = document.getElementById('excelToPdf').files[0];
                if (!file) return alert('Selecione um arquivo Excel');
                document.getElementById('loadingExcelToPdf').style.display = 'block';
                document.getElementById('resultExcelToPdf').style.display = 'none';
                const formData = new FormData();
                formData.append('file', file);
                try {
                    const response = await fetch('/convert/excel-to-pdf', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultExcelToPdf').innerHTML = `<a href="${url}" download="${file.name.replace(/\.[^.]+$/, '.pdf')}">📥 Download PDF</a>`;
                        document.getElementById('resultExcelToPdf').className = 'result';
                    } else {
                        throw new Error('Erro na conversão');
                    }
                } catch (error) {
                    document.getElementById('resultExcelToPdf').innerHTML = 'Erro ao converter';
                    document.getElementById('resultExcelToPdf').className = 'result error';
                } finally {
                    document.getElementById('loadingExcelToPdf').style.display = 'none';
                    document.getElementById('resultExcelToPdf').style.display = 'block';
                }
            }
            async function convertPptToPdf() {
                const file = document.getElementById('pptToPdf').files[0];
                if (!file) return alert('Selecione um arquivo PowerPoint');
                document.getElementById('loadingPptToPdf').style.display = 'block';
                document.getElementById('resultPptToPdf').style.display = 'none';
                const formData = new FormData();
                formData.append('file', file);
                try {
                    const response = await fetch('/convert/ppt-to-pdf', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultPptToPdf').innerHTML = `<a href="${url}" download="${file.name.replace(/\.[^.]+$/, '.pdf')}">📥 Download PDF</a>`;
                        document.getElementById('resultPptToPdf').className = 'result';
                    } else {
                        throw new Error('Erro na conversão');
                    }
                } catch (error) {
                    document.getElementById('resultPptToPdf').innerHTML = 'Erro ao converter';
                    document.getElementById('resultPptToPdf').className = 'result error';
                } finally {
                    document.getElementById('loadingPptToPdf').style.display = 'none';
                    document.getElementById('resultPptToPdf').style.display = 'block';
                }
            }
            async function convertHtmlToPdf() {
                const file = document.getElementById('htmlToPdf').files[0];
                if (!file) return alert('Selecione um arquivo HTML');
                document.getElementById('loadingHtmlToPdf').style.display = 'block';
                document.getElementById('resultHtmlToPdf').style.display = 'none';
                const formData = new FormData();
                formData.append('file', file);
                try {
                    const response = await fetch('/convert/html-to-pdf', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultHtmlToPdf').innerHTML = `<a href="${url}" download="${file.name.replace(/\.[^.]+$/, '.pdf')}">📥 Download PDF</a>`;
                        document.getElementById('resultHtmlToPdf').className = 'result';
                    } else {
                        throw new Error('Erro na conversão');
                    }
                } catch (error) {
                    document.getElementById('resultHtmlToPdf').innerHTML = 'Erro ao converter';
                    document.getElementById('resultHtmlToPdf').className = 'result error';
                } finally {
                    document.getElementById('loadingHtmlToPdf').style.display = 'none';
                    document.getElementById('resultHtmlToPdf').style.display = 'block';
                }
            }

            async function unlockPdf() {
                const file = document.getElementById('unlockPdf').files[0];
                const password = document.getElementById('unlockPassword').value;
                if (!file || !password) return alert('Selecione o PDF e informe a senha');
                document.getElementById('loadingUnlockPdf').style.display = 'block';
                document.getElementById('resultUnlockPdf').style.display = 'none';
                const formData = new FormData();
                formData.append('file', file);
                formData.append('password', password);
                try {
                    const response = await fetch('/unlock-pdf', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultUnlockPdf').innerHTML = `<a href="${url}" download="unlocked.pdf">📥 Download PDF Desbloqueado</a>`;
                        document.getElementById('resultUnlockPdf').className = 'result';
                    } else {
                        throw new Error('Erro ao desbloquear');
                    }
                } catch (error) {
                    document.getElementById('resultUnlockPdf').innerHTML = 'Erro ao desbloquear PDF';
                    document.getElementById('resultUnlockPdf').className = 'result error';
                } finally {
                    document.getElementById('loadingUnlockPdf').style.display = 'none';
                    document.getElementById('resultUnlockPdf').style.display = 'block';
                }
            }
            async function protectPdf() {
                const file = document.getElementById('protectPdf').files[0];
                const password = document.getElementById('protectPassword').value;
                if (!file || !password) return alert('Selecione o PDF e informe a nova senha');
                document.getElementById('loadingProtectPdf').style.display = 'block';
                document.getElementById('resultProtectPdf').style.display = 'none';
                const formData = new FormData();
                formData.append('file', file);
                formData.append('password', password);
                try {
                    const response = await fetch('/protect-pdf', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultProtectPdf').innerHTML = `<a href="${url}" download="protected.pdf">📥 Download PDF Protegido</a>`;
                        document.getElementById('resultProtectPdf').className = 'result';
                    } else {
                        throw new Error('Erro ao proteger');
                    }
                } catch (error) {
                    document.getElementById('resultProtectPdf').innerHTML = 'Erro ao proteger PDF';
                    document.getElementById('resultProtectPdf').className = 'result error';
                } finally {
                    document.getElementById('loadingProtectPdf').style.display = 'none';
                    document.getElementById('resultProtectPdf').style.display = 'block';
                }
            }

            async function insertPageNumbers() {
                const file = document.getElementById('numberPdf').files[0];
                if (!file) return alert('Selecione um PDF');
                document.getElementById('loadingNumberPdf').style.display = 'block';
                document.getElementById('resultNumberPdf').style.display = 'none';
                const formData = new FormData();
                formData.append('file', file);
                try {
                    const response = await fetch('/edit/add-page-numbers', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultNumberPdf').innerHTML = `<a href="${url}" download="numbered.pdf">📥 Download PDF Numerado</a>`;
                        document.getElementById('resultNumberPdf').className = 'result';
                    } else {
                        throw new Error('Erro ao inserir números');
                    }
                } catch (error) {
                    document.getElementById('resultNumberPdf').innerHTML = 'Erro ao inserir números';
                    document.getElementById('resultNumberPdf').className = 'result error';
                } finally {
                    document.getElementById('loadingNumberPdf').style.display = 'none';
                    document.getElementById('resultNumberPdf').style.display = 'block';
                }
            }
            async function insertWatermark() {
                const file = document.getElementById('watermarkPdf').files[0];
                const text = document.getElementById('watermarkText').value;
                if (!file || !text) return alert('Selecione o PDF e informe o texto da marca d\'água');
                document.getElementById('loadingWatermarkPdf').style.display = 'block';
                document.getElementById('resultWatermarkPdf').style.display = 'none';
                const formData = new FormData();
                formData.append('file', file);
                formData.append('text', text);
                try {
                    const response = await fetch('/edit/add-watermark', { method: 'POST', body: formData });
                    if (response.ok) {
                        const blob = await response.blob();
                        const url = URL.createObjectURL(blob);
                        document.getElementById('resultWatermarkPdf').innerHTML = `<a href="${url}" download="watermarked.pdf">📥 Download PDF com Marca d'Água</a>`;
                        document.getElementById('resultWatermarkPdf').className = 'result';
                    } else {
                        throw new Error('Erro ao inserir marca d\'água');
                    }
                } catch (error) {
                    document.getElementById('resultWatermarkPdf').innerHTML = 'Erro ao inserir marca d\'água';
                    document.getElementById('resultWatermarkPdf').className = 'result error';
                } finally {
                    document.getElementById('loadingWatermarkPdf').style.display = 'none';
                    document.getElementById('resultWatermarkPdf').style.display = 'block';
                }
            }

            // Alternar categorias na barra lateral
            document.querySelectorAll('.category-link').forEach(link => {
                link.addEventListener('click', function(e) {
                    e.preventDefault();
                    document.querySelectorAll('.category-link').forEach(l => l.classList.remove('active'));
                    this.classList.add('active');
                    // Esconde todas as grids
                    document.querySelectorAll('.tools-grid').forEach(grid => grid.style.display = 'none');
                    // Mostra a grid da categoria
                    const cat = this.getAttribute('data-category');
                    document.getElementById(cat).style.display = 'grid';
                    // Atualiza título
                    document.getElementById('category-title').textContent = this.textContent;
                    // Atualiza descrição
                    let desc = '';
                    switch(cat) {
                        case 'organizar': desc = 'Manipule e organize seus arquivos PDF.'; break;
                        case 'otimizar': desc = 'Otimize e melhore seus PDFs.'; break;
                        case 'converter-em': desc = 'Converta outros formatos em PDF.'; break;
                        case 'converter-de': desc = 'Converta PDF para outros formatos.'; break;
                        case 'editar': desc = 'Edite e personalize seus PDFs.'; break;
                        case 'seguranca': desc = 'Proteja e gerencie a segurança dos seus PDFs.'; break;
                        default: desc = 'Manipule seus arquivos PDF facilmente';
                    }
                    document.getElementById('category-desc').textContent = desc;
                });
            });
        </script>
        <script src="/static/app.js"></script>
    </body>
    </html>
    """

@app.post("/convert/word")
async def convert_to_word(file: UploadFile = File(...)):
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Arquivo deve ser PDF")
    
    # Salvar arquivo temporário
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        # Extrair texto do PDF
        with open(temp_path, 'rb') as pdf_file:
            reader = pypdf.PdfReader(pdf_file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
        
        # Criar documento Word
        doc = Document()
        doc.add_paragraph(text)
        
        # Salvar e retornar
        output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_converted.docx"
        doc.save(output_path)
        
        return FileResponse(output_path, filename=file.filename.replace('.pdf', '.docx'))
    
    finally:
        # Limpar arquivo temporário
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/convert/excel")
async def convert_to_excel(file: UploadFile = File(...)):
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Arquivo deve ser PDF")
    
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        # Extrair texto do PDF
        with open(temp_path, 'rb') as pdf_file:
            reader = pypdf.PdfReader(pdf_file)
            
        # Criar planilha Excel
        wb = Workbook()
        ws = wb.active
        ws.title = "PDF Content"
        
        row = 1
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            ws.cell(row=row, column=1, value=f"Página {page_num}")
            ws.cell(row=row, column=2, value=text)
            row += 1
        
        output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_converted.xlsx"
        wb.save(output_path)
        
        return FileResponse(output_path, filename=file.filename.replace('.pdf', '.xlsx'))
    
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/merge")
async def merge_pdfs(files: List[UploadFile] = File(...)):
    if len(files) < 2:
        raise HTTPException(status_code=400, detail="Necessário pelo menos 2 arquivos")
    
    temp_paths = []
    try:
        # Salvar arquivos temporários
        for file in files:
            if not file.filename.endswith('.pdf'):
                raise HTTPException(status_code=400, detail="Todos os arquivos devem ser PDF")
            
            temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
            with open(temp_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)
            temp_paths.append(temp_path)
        
        # Juntar PDFs
        merger = pypdf.PdfWriter()
        for path in temp_paths:
            merger.append(path)
        
        output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_merged.pdf"
        with open(output_path, 'wb') as output_file:
            merger.write(output_file)
        
        return FileResponse(output_path, filename="merged.pdf")
    
    finally:
        # Limpar arquivos temporários
        for path in temp_paths:
            if os.path.exists(path):
                os.remove(path)

@app.post("/split")
async def split_pdf(file: UploadFile = File(...), start_page: int = Form(...), end_page: int = Form(...)):
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Arquivo deve ser PDF")
    
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        with open(temp_path, 'rb') as pdf_file:
            reader = pypdf.PdfReader(pdf_file)
            
            if start_page < 1 or end_page > len(reader.pages) or start_page > end_page:
                raise HTTPException(status_code=400, detail="Páginas inválidas")
            
            writer = pypdf.PdfWriter()
            for i in range(start_page - 1, end_page):
                writer.add_page(reader.pages[i])
            
            output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_extracted.pdf"
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            return FileResponse(output_path, filename="extracted_pages.pdf")
    
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/compress")
async def compress_pdf(file: UploadFile = File(...)):
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Arquivo deve ser PDF")
    
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        with open(temp_path, 'rb') as pdf_file:
            reader = pypdf.PdfReader(pdf_file)
            writer = pypdf.PdfWriter()
            
            for page in reader.pages:
                writer.add_page(page)
            
            # Compressão básica
            writer.compress_identical_objects()
            
            output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_compressed.pdf"
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            return FileResponse(output_path, filename=file.filename.replace('.pdf', '_compressed.pdf'))
    
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/compare")
async def compare_pdfs(file1: UploadFile = File(...), file2: UploadFile = File(...)):
    temp_paths = []
    try:
        # Salvar arquivos temporários
        for file in [file1, file2]:
            if not file.filename.endswith('.pdf'):
                raise HTTPException(status_code=400, detail="Arquivos devem ser PDF")
            
            temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
            with open(temp_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)
            temp_paths.append(temp_path)
        
        # Extrair texto dos PDFs
        texts = []
        page_counts = []
        
        for path in temp_paths:
            with open(path, 'rb') as pdf_file:
                reader = pypdf.PdfReader(pdf_file)
                page_counts.append(len(reader.pages))
                text = ""
                for page in reader.pages:
                    text += page.extract_text()
                texts.append(text)
        
        # Comparação básica
        similarity = len(set(texts[0].split()) & set(texts[1].split())) / len(set(texts[0].split()) | set(texts[1].split())) * 100
        
        return {
            "message": f"📊 Comparação concluída:<br>• Arquivo 1: {page_counts[0]} páginas<br>• Arquivo 2: {page_counts[1]} páginas<br>• Similaridade: {similarity:.1f}%"
        }
    
    finally:
        for path in temp_paths:
            if os.path.exists(path):
                os.remove(path)

# NOVOS ENDPOINTS FASTAPI
@app.post("/remove-pages")
async def remove_pages(file: UploadFile = File(...), pages: str = Form(...)):
    import re
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    try:
        with open(temp_path, 'rb') as pdf_file:
            reader = pypdf.PdfReader(pdf_file)
            writer = pypdf.PdfWriter()
            total = len(reader.pages)
            # Parse páginas a remover
            remove_set = set()
            for part in pages.split(','):
                if '-' in part:
                    start, end = map(int, part.split('-'))
                    remove_set.update(range(start, end+1))
                else:
                    remove_set.add(int(part))
            for i in range(total):
                if (i+1) not in remove_set:
                    writer.add_page(reader.pages[i])
            output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_removed.pdf"
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            return FileResponse(output_path, filename="removed_pages.pdf")
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/extract-pages")
async def extract_pages(file: UploadFile = File(...), pages: str = Form(...)):
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    try:
        with open(temp_path, 'rb') as pdf_file:
            reader = pypdf.PdfReader(pdf_file)
            writer = pypdf.PdfWriter()
            total = len(reader.pages)
            extract_set = set()
            for part in pages.split(','):
                if '-' in part:
                    start, end = map(int, part.split('-'))
                    extract_set.update(range(start, end+1))
                else:
                    extract_set.add(int(part))
            for i in range(total):
                if (i+1) in extract_set:
                    writer.add_page(reader.pages[i])
            output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_extracted.pdf"
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            return FileResponse(output_path, filename="extracted_pages.pdf")
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/organize-pages")
async def organize_pages(file: UploadFile = File(...), order: str = Form(...)):
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    try:
        with open(temp_path, 'rb') as pdf_file:
            reader = pypdf.PdfReader(pdf_file)
            writer = pypdf.PdfWriter()
            total = len(reader.pages)
            order_list = [int(x) for x in order.split(',') if x.strip().isdigit()]
            for idx in order_list:
                if 1 <= idx <= total:
                    writer.add_page(reader.pages[idx-1])
            output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_organized.pdf"
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            return FileResponse(output_path, filename="organized.pdf")
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

# NOVOS ENDPOINTS CONVERTER EM PDF
from fastapi import Request
from fastapi.responses import StreamingResponse
from PIL import Image
import io

@app.post("/convert/jpg-to-pdf")
async def jpg_to_pdf(files: list[UploadFile] = File(...)):
    images = []
    temp_paths = []
    for file in files:
        temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
        with open(temp_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        temp_paths.append(temp_path)
        img = Image.open(temp_path).convert("RGB")
        images.append(img)
    output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_jpg2pdf.pdf"
    images[0].save(output_path, save_all=True, append_images=images[1:])
    for path in temp_paths:
        if os.path.exists(path): os.remove(path)
    return FileResponse(output_path, filename="imagens.pdf")

@app.post("/convert/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    from docx import Document
    from fpdf import FPDF
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    doc = Document(temp_path)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)
    for para in doc.paragraphs:
        pdf.multi_cell(0, 10, para.text)
    output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_word2pdf.pdf"
    pdf.output(output_path)
    if os.path.exists(temp_path): os.remove(temp_path)
    return FileResponse(output_path, filename=file.filename.replace('.docx', '.pdf').replace('.doc', '.pdf'))

@app.post("/convert/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    import pandas as pd
    from fpdf import FPDF
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    df = pd.read_excel(temp_path)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=10)
    col_width = pdf.w / (len(df.columns) + 1)
    row_height = pdf.font_size * 1.5
    # Cabeçalho
    for col in df.columns:
        pdf.cell(col_width, row_height, str(col), border=1)
    pdf.ln(row_height)
    # Linhas
    for i, row in df.iterrows():
        for item in row:
            pdf.cell(col_width, row_height, str(item), border=1)
        pdf.ln(row_height)
    output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_excel2pdf.pdf"
    pdf.output(output_path)
    if os.path.exists(temp_path): os.remove(temp_path)
    return FileResponse(output_path, filename=file.filename.replace('.xlsx', '.pdf').replace('.xls', '.pdf'))

@app.post("/convert/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    from pptx import Presentation
    from fpdf import FPDF
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    prs = Presentation(temp_path)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=14)
    for slide in prs.slides:
        pdf.add_page()
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                pdf.multi_cell(0, 10, shape.text.strip())
    output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_ppt2pdf.pdf"
    pdf.output(output_path)
    if os.path.exists(temp_path): os.remove(temp_path)
    return FileResponse(output_path, filename=file.filename.replace('.pptx', '.pdf').replace('.ppt', '.pdf'))

@app.post("/convert/html-to-pdf")
async def html_to_pdf(file: UploadFile = File(...)):
    from bs4 import BeautifulSoup
    from fpdf import FPDF
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    with open(temp_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'html.parser')
        text = soup.get_text()
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, text)
    output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_html2pdf.pdf"
    pdf.output(output_path)
    if os.path.exists(temp_path): os.remove(temp_path)
    return FileResponse(output_path, filename=file.filename.replace('.html', '.pdf').replace('.htm', '.pdf'))

# ENDPOINTS SEGURANÇA DO PDF
@app.post("/unlock-pdf")
async def unlock_pdf(file: UploadFile = File(...), password: str = Form(...)):
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    try:
        with open(temp_path, 'rb') as pdf_file:
            reader = pypdf.PdfReader(pdf_file, password=password)
            writer = pypdf.PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_unlocked.pdf"
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            return FileResponse(output_path, filename="unlocked.pdf")
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/protect-pdf")
async def protect_pdf(file: UploadFile = File(...), password: str = Form(...)):
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    try:
        with open(temp_path, 'rb') as pdf_file:
            reader = pypdf.PdfReader(pdf_file)
            writer = pypdf.PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            writer.encrypt(password)
            output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_protected.pdf"
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            return FileResponse(output_path, filename="protected.pdf")
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

# ENDPOINTS EDITAR PDF
@app.post("/edit/add-page-numbers")
async def add_page_numbers(file: UploadFile = File(...)):
    from PyPDF2 import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    import io
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_numbered.pdf"
    try:
        reader = PdfReader(temp_path)
        writer = PdfWriter()
        for i, page in enumerate(reader.pages):
            packet = io.BytesIO()
            can = canvas.Canvas(packet, pagesize=letter)
            can.setFont("Helvetica", 10)
            can.drawString(500, 20, f"{i+1}")
            can.save()
            packet.seek(0)
            from PyPDF2 import PdfReader as PR
            watermark = PR(packet)
            page.merge_page(watermark.pages[0])
            writer.add_page(page)
        with open(output_path, 'wb') as f:
            writer.write(f)
        return FileResponse(output_path, filename="numbered.pdf")
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/edit/add-watermark")
async def add_watermark(file: UploadFile = File(...), text: str = Form(...)):
    from PyPDF2 import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    import io
    temp_path = f"{UPLOAD_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    output_path = f"{OUTPUT_DIR}/{uuid.uuid4()}_watermarked.pdf"
    try:
        reader = PdfReader(temp_path)
        writer = PdfWriter()
        for page in reader.pages:
            packet = io.BytesIO()
            can = canvas.Canvas(packet, pagesize=letter)
            can.setFont("Helvetica", 16)
            can.setFillColorRGB(0.7, 0.7, 0.7)
            can.saveState()
            can.translate(300, 400)
            can.rotate(30)
            can.drawCentredString(0, 0, text)
            can.restoreState()
            can.save()
            packet.seek(0)
            from PyPDF2 import PdfReader as PR
            watermark = PR(packet)
            page.merge_page(watermark.pages[0])
            writer.add_page(page)
        with open(output_path, 'wb') as f:
            writer.write(f)
        return FileResponse(output_path, filename="watermarked.pdf")
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run("main:app", host="0.0.0.0", port=port, reload=False)